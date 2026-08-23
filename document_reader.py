"""
document_reader.py — extract text from PDFs, Word docs, PowerPoint files, and images.
All free, all local, nothing uploaded anywhere.
"""

import os
import pdfplumber
from docx import Document
from pptx import Presentation
from PIL import Image
import pytesseract


def read_pdf(path):
    text = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text.append(page_text)
    return "\n".join(text)


def read_docx(path):
    doc = Document(path)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def read_pptx(path):
    prs = Presentation(path)
    text = []
    for i, slide in enumerate(prs.slides, 1):
        slide_text = [f"--- Slide {i} ---"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text)
        text.append("\n".join(slide_text))
    return "\n".join(text)


def read_image(path):
    """OCR: extract any text that appears inside an image."""
    img = Image.open(path)
    return pytesseract.image_to_string(img)


def read_document(path):
    """Auto-detect file type and extract text accordingly."""
    if not os.path.exists(path):
        return f"File not found: {path}"

    ext = os.path.splitext(path)[1].lower()
    try:
        if ext == ".pdf":
            return read_pdf(path)
        elif ext == ".docx":
            return read_docx(path)
        elif ext == ".pptx":
            return read_pptx(path)
        elif ext in (".png", ".jpg", ".jpeg", ".webp", ".bmp"):
            return read_image(path)
        elif ext == ".txt":
            with open(path, "r", errors="ignore") as f:
                return f.read()
        else:
            return f"Unsupported file type: {ext}"
    except Exception as e:
        return f"Error reading {path}: {e}"


if __name__ == "__main__":
    # Quick manual test: python3 document_reader.py documents/somefile.pdf
    import sys
    if len(sys.argv) > 1:
        print(read_document(sys.argv[1]))
